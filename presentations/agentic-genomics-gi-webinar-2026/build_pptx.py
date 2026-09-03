#!/usr/bin/env python3
"""Export the ClawBio x GI webinar HTML deck to an editable PPTX.

Text is authored as native PowerPoint text boxes and shapes (fully editable).
The four bespoke SVG diagrams are embedded as high-resolution PNGs (rendered
from the inline SVGs via rsvg-convert in the build step) because vector
diagrams cannot be turned into native PowerPoint shapes cleanly.

Run from this directory:
    /opt/homebrew/bin/python3 build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
SVG = os.path.join(HERE, "build", "svg")
OUT = os.path.join(HERE, "agentic-genomics-gi-webinar-2026.pptx")

# palette (matches the deck CSS)
BG       = RGBColor(0x0d, 0x11, 0x17)
TEXT     = RGBColor(0xe6, 0xed, 0xf3)
MUTED    = RGBColor(0xad, 0xba, 0xc7)
DIM      = RGBColor(0x6e, 0x76, 0x81)
ACCENT   = RGBColor(0x58, 0xa6, 0xff)
GREEN    = RGBColor(0x56, 0xd3, 0x64)
ORANGE   = RGBColor(0xe3, 0xb3, 0x41)
RED      = RGBColor(0xff, 0xa1, 0x98)
PURPLE   = RGBColor(0xd2, 0xa8, 0xff)
CARDBG   = RGBColor(0x16, 0x1b, 0x22)
BORDER   = RGBColor(0x30, 0x36, 0x3d)
FONT     = "Arial"

SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _set_runs(p, runs, size, align=PP_ALIGN.CENTER, bold=False):
    p.alignment = align
    for item in runs:
        txt, color = item[0], item[1]
        b = item[2] if len(item) > 2 else bold
        r = p.add_run(); r.text = txt
        f = r.font
        f.size = Pt(size); f.name = FONT; f.bold = b; f.color.rgb = color


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {runs,size,align,bold,space_after}."""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_runs(p, ln["runs"], ln.get("size", 16),
                  ln.get("align", PP_ALIGN.CENTER), ln.get("bold", False))
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        if "space_before" in ln:
            p.space_before = Pt(ln["space_before"])
    return box


def chip(s, text, color=ACCENT, top=0.45, cx=SW / 2):
    w = max(1.4, len(text) * 0.085 + 0.5); h = 0.34
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(cx - w / 2), Inches(top), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    _set_runs(p, [(text.upper(), BG, True)], 11, PP_ALIGN.CENTER)
    return shp


def title(s, runs, top=0.92, size=30):
    return textbox(s, 0.6, top, SW - 1.2, 1.0,
                   [{"runs": runs, "size": size, "bold": True, "align": PP_ALIGN.CENTER}])


def card(s, l, t, w, h, head, head_color, body):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    shp.adjustments[0] = 0.06
    shp.fill.solid(); shp.fill.fore_color.rgb = CARDBG
    shp.line.color.rgb = BORDER; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    _set_runs(p, [(head, head_color, True)], 15, PP_ALIGN.CENTER); p.space_after = Pt(6)
    p2 = tf.add_paragraph()
    _set_runs(p2, [(body, MUTED)], 12.5, PP_ALIGN.LEFT)
    return shp


def bullets(s, l, t, w, h, head, head_color, items, head_size=17, size=14.5):
    lines = []
    if head:
        lines.append({"runs": [(head, head_color, True)], "size": head_size,
                      "align": PP_ALIGN.LEFT, "space_after": 8})
    for it in items:
        txt, color = (it if isinstance(it, tuple) else (it, TEXT))
        lines.append({"runs": [("→  ", ACCENT), (txt, color)], "size": size,
                      "align": PP_ALIGN.LEFT, "space_after": 6})
    return textbox(s, l, t, w, h, lines)


def stat(s, l, t, w, num, ncolor, label, sub=None, num_size=30):
    lines = [{"runs": [(num, ncolor, True)], "size": num_size, "align": PP_ALIGN.CENTER,
              "space_after": 2}]
    lines.append({"runs": [(label, MUTED)], "size": 12, "align": PP_ALIGN.CENTER})
    if sub:
        lines.append({"runs": [(sub, DIM)], "size": 10.5, "align": PP_ALIGN.CENTER})
    return textbox(s, l, t, w, 1.3, lines)


def picture(s, path, max_w, max_h, cx=None, top=None, left=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    if left is None:
        left = (cx if cx is not None else SW / 2) - w / 2
    return s.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))


def caption(s, top, runs, size=17):
    return textbox(s, 0.8, top, SW - 1.6, 0.9,
                   [{"runs": runs, "size": size, "align": PP_ALIGN.CENTER}])


# ---------------------------------------------------------------- slide 1: title
s = slide()
picture(s, os.path.join(FIG, "clawbio-logo-white.png"), 1.4, 1.05, cx=SW/2, top=0.55)
chip(s, "ClawBio × Genomic Intelligence", ACCENT, top=1.78)
textbox(s, 0.6, 2.25, SW - 1.2, 1.2,
        [{"runs": [("Agentic Genomics in ", TEXT, True), ("Practice", GREEN, True)],
          "size": 42, "align": PP_ALIGN.CENTER}])
textbox(s, 1.8, 3.55, SW - 3.6, 0.9,
        [{"runs": [("AI agents are starting to do real genomics work end to end. "
                    "The hard part is no longer writing the code. It is trusting the result.",
                    MUTED)], "size": 16, "align": PP_ALIGN.CENTER}])
textbox(s, 0.6, 4.75, SW - 1.2, 1.6, [
    {"runs": [("Manuel Corpas", ACCENT, True)], "size": 19, "align": PP_ALIGN.CENTER, "space_after": 6},
    {"runs": [("Senior Lecturer in Genomics, AI, and Data Science  ·  University of Westminster", DIM)], "size": 12.5, "align": PP_ALIGN.CENTER},
    {"runs": [("Founder, ClawBio", DIM)], "size": 12.5, "align": PP_ALIGN.CENTER},
    {"runs": [("Webinar  ·  ClawBio × Genomic Intelligence  ·  24 June 2026", DIM)], "size": 12.5, "align": PP_ALIGN.CENTER},
])

# ---------------------------------------------------------------- slide 2: agenda
s = slide()
chip(s, "Today")
title(s, [("What we will cover in the next 90 minutes", TEXT)], size=28)
rows = [
    ("Manuel (ClawBio)", ACCENT, "State of the art in agentic genomics: what changed, why trust is the bottleneck, and concrete examples with ClawBio skills."),
    ("Benjamin (GI)", PURPLE, "What the Genomic Intelligence models are and where they fit in a genomics workflow."),
    ("Benjamin (GI)", PURPLE, "Live scenarios: using GI skills inside ClawBio to solve real bioinformatics problems, in Codex and on the web platform."),
    ("Everyone", ACCENT, "Q&A, and an open call: build your own use case with the platform over the next two weeks."),
]
ry = 2.05
for who, wc, what in rows:
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(ry), Inches(11.13), Inches(0.92))
    shp.adjustments[0] = 0.12
    shp.fill.solid(); shp.fill.fore_color.rgb = CARDBG
    shp.line.color.rgb = BORDER; shp.line.width = Pt(1); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    _set_runs(p, [(who + "   ", wc, True), (what, MUTED)], 13.5, PP_ALIGN.LEFT)
    ry += 1.04
textbox(s, 1.1, ry + 0.05, 11.13, 0.5,
        [{"runs": [("For bioinformaticians, computational biologists, genomics researchers, "
                    "and anyone building with AI agents in the life sciences.", DIM)],
          "size": 11.5, "align": PP_ALIGN.CENTER}])

# ---------------------------------------------------------------- slide 3: two waves
s = slide()
chip(s, "The shift")
title(s, [("Two Waves of LLMs in Biology", TEXT)], size=28)
picture(s, os.path.join(SVG, "diagram-0.png"), 11.6, 1.5, cx=SW/2, top=1.95)
bullets(s, 0.9, 3.7, 5.7, 2.6, "First wave: information retrieval", MUTED, [
    "Summarising papers, answering pathway questions, extracting structured data from text.",
    ("Useful but incremental.", DIM),
])
bullets(s, 6.85, 3.7, 5.7, 2.6, "Second wave: autonomous execution", GREEN, [
    "Modern LLMs can write, debug, and execute code.",
    "Connected to file systems, databases, and command-line tools, they plan multi-step analyses and adapt on intermediate results.",
])

# ---------------------------------------------------------------- slide 4: defining
s = slide()
chip(s, "Definition")
title(s, [("Defining ", TEXT), ("Agentic Genomics", GREEN)], size=28)
textbox(s, 1.4, 1.85, SW - 2.8, 1.1,
        [{"runs": [("An AI agent that plans and executes a genomic analysis end to end: it reads the data, "
                    "chooses the tools, runs them, inspects intermediate results, and revises its own plan, "
                    "with a human evaluating rather than typing every command.", TEXT)],
          "size": 16, "align": PP_ALIGN.CENTER}])
cw, gap = 3.7, 0.45
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (h, b) in enumerate([
    ("Senses", "Reads VCFs, FASTQs, manifests, guidelines, and prior outputs as context."),
    ("Acts", "Invokes real tools: aligners, variant callers, annotation and reporting skills."),
    ("Adapts", "Checks intermediate results and re-plans, instead of running a fixed script."),
]):
    card(s, x0 + i * (cw + gap), 3.2, cw, 1.7, h, ACCENT, b)
textbox(s, 1.0, 5.25, SW - 2.0, 0.7,
        [{"runs": [("Grounded in the Corpas, Fatumo & Guio Perspective on agentic genomics, and a Briefings "
                    "in Bioinformatics benchmark of frontier LLMs in pharmacogenomics.", DIM)],
          "size": 11.5, "align": PP_ALIGN.CENTER}])

# ---------------------------------------------------------------- slide 5: paradigm shift
s = slide()
chip(s, "The paradigm shift")
title(s, [("The work moves from ", TEXT), ("code production", ACCENT), (" to ", TEXT), ("validation", GREEN)], size=27)
picture(s, os.path.join(SVG, "diagram-1.png"), 11.8, 4.0, cx=SW/2, top=1.95)
caption(s, 6.5, [("Generating the analysis is no longer the hard part. ", TEXT),
                 ("Knowing you can trust it", GREEN), (" is.", TEXT)])

# ---------------------------------------------------------------- slide 6: bottleneck
s = slide()
chip(s, "The new bottleneck")
title(s, [("The Validation Bottleneck: ", TEXT), ("silent, plausible-looking failure", RED)], size=26)
cw, gap = 3.7, 0.45
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (h, hc, b) in enumerate([
    ("Plausible", RED, "The output is formatted correctly, reads fluently, and matches what a reviewer expects to see."),
    ("Confident", ORANGE, "The agent states a wrong answer with the same tone as a right one. No flag, no uncertainty."),
    ("Consequential", ACCENT, "In genomics the failure mode is a missed cancer gene or a wrong dose, not a 404. And it scales."),
]):
    card(s, x0 + i * (cw + gap), 1.85, cw, 1.6, h, hc, b)
# dramatic box
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(3.7), Inches(11.13), Inches(1.85))
box.adjustments[0] = 0.05
box.fill.solid(); box.fill.fore_color.rgb = CARDBG
box.line.color.rgb = RED; box.line.width = Pt(1.5); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
p = tf.paragraphs[0]
_set_runs(p, [
    ("A healthy 32-year-old's exome is screened. The agent calls a pathogenic ", TEXT),
    ("BRCA1", RED), (" frameshift ", TEXT), ("“likely benign”", MUTED),
    (", with a clean report, a citation, and ", TEXT), ("no flag raised", RED),
    (". She is never offered screening or risk-reducing surgery. The first sign is ", TEXT),
    ("stage IV cancer", RED), (".", TEXT),
], 14.5, PP_ALIGN.LEFT)
p.space_after = Pt(8)
p2 = tf.add_paragraph()
_set_runs(p2, [("Nothing crashed. Now run that same pipeline across ", TEXT),
               ("10,000 genomes overnight", RED), (".", TEXT)], 14.5, PP_ALIGN.LEFT)
caption(s, 5.75, [("A confident, plausible, wrong answer is worse than no answer. ", TEXT),
                  ("Safe uncertainty beats confident hallucination.", GREEN)], size=15)

# ---------------------------------------------------------------- slide 7: ClawBio
s = slide()
textbox(s, 0.6, 0.45, SW - 1.2, 0.9,
        [{"runs": [("Claw", TEXT, True), ("Bio", GREEN, True)], "size": 40, "align": PP_ALIGN.CENTER}])
textbox(s, 0.6, 1.35, SW - 1.2, 0.7, [
    {"runs": [("An open, agent-native skill library for bioinformatics.", MUTED)], "size": 16, "align": PP_ALIGN.CENTER},
    {"runs": [("Open-source  ·  local-first  ·  reproducible", DIM)], "size": 12, "align": PP_ALIGN.CENTER},
])
sw_, sgap = 2.6, 0.3
sx0 = (SW - (sw_ * 4 + sgap * 3)) / 2
for i, (n, l) in enumerate([("993", "GitHub stars"), ("90", "open skills"), ("48", "contributors"), ("218", "forks")]):
    stat(s, sx0 + i * (sw_ + sgap), 2.35, sw_, n, GREEN, l, num_size=32)
textbox(s, 1.0, 3.35, SW - 2.0, 0.4,
        [{"runs": [("From a standing start in February 2026. Nearly ", DIM),
                   ("1,000 stars in four months", GREEN), (", heading for 5,000.", DIM)],
          "size": 12, "align": PP_ALIGN.CENTER}])
cw, gap = 3.7, 0.45
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (h, b) in enumerate([
    ("A skill is a contract", "Each skill is a plain-text SKILL.md plus code: what it does, what it needs, and the steps the agent must follow."),
    ("The connecting layer", "ClawBio links AI agents to genomic databases, clinical knowledge, and validation workflows."),
    ("Open and shared", "Anyone can write, audit, and reuse skills. The aim: a trust layer for agentic genomics."),
]):
    card(s, x0 + i * (cw + gap), 3.95, cw, 1.85, h, ACCENT, b)
textbox(s, 1.0, 6.0, SW - 2.0, 0.5,
        [{"runs": [("github.com/ClawBio/ClawBio  ·  ", GREEN),
                   ("today, Genomic Intelligence's models run as skills inside this library.", DIM)],
          "size": 12, "align": PP_ALIGN.CENTER}])

# ---------------------------------------------------------------- slide 8: empirical
s = slide()
chip(s, "Empirical question")
title(s, [("Can a plain-text SKILL.md reach ", TEXT), ("clinical-grade", GREEN), ("?", TEXT)], size=27)
textbox(s, 1.4, 1.9, SW - 2.8, 0.9,
        [{"runs": [("Pharmacogenomics, ground truth from CPIC guidelines, ", TEXT),
                   ("44,550 scored evaluations", GREEN),
                   (" across frontier models and three conditions: free prompt, the skill's reasoning, "
                    "and executing the skill.", TEXT)], "size": 15, "align": PP_ALIGN.CENTER}])
sw_, sgap = 3.4, 0.5
sx0 = (SW - (sw_ * 3 + sgap * 2)) / 2
for i, (n, c, l, sub) in enumerate([
    ("80.6%", RED, "Free prompt", "(no specification)"),
    ("95.5%", ORANGE, "Skill reasoning", "(guideline loaded)"),
    ("100%", GREEN, "Skill execution", "(deterministic)"),
]):
    stat(s, sx0 + i * (sw_ + sgap), 3.3, sw_, n, c, l, sub, num_size=40)
caption(s, 5.4, [("Accuracy climbs as correctness is ", TEXT),
                 ("constrained by architecture", GREEN),
                 (", not by asking the model more nicely.", TEXT)])

# ---------------------------------------------------------------- slide 9: counterintuitive
s = slide()
chip(s, "The counterintuitive result")
title(s, [("Giving the model the right guideline made it ", TEXT), ("more dangerous", RED)], size=25)
picture(s, os.path.join(SVG, "diagram-2.png"), 11.6, 3.7, cx=SW/2, top=2.0)
caption(s, 6.2, [("Trust is ", TEXT), ("architectural", GREEN),
                 (": it comes from deterministic, auditable, model-invariant execution, not from a better prompt.", TEXT)], size=15)

# ---------------------------------------------------------------- slide 10: ancestry
s = slide()
chip(s, "And it is not equal")
title(s, [("On real genomes, accuracy ", TEXT), ("falls by ancestry", RED)], size=27)
picture(s, os.path.join(SVG, "diagram-3.png"), 11.0, 3.9, cx=SW/2, top=1.95)
caption(s, 6.25, [("Curated ~96% ", TEXT), ("does not transfer", RED),
                  (" to real diplotypes from over 7,000 individuals. ", TEXT),
                  ("Executing the skill removes the gradient.", GREEN),
                  (" Validation is also an equity problem.", TEXT)], size=14)

# ---------------------------------------------------------------- slide 11: handoff
s = slide()
chip(s, "Over to Genomic Intelligence", PURPLE)
title(s, [("From the framing to the ", TEXT), ("live work", GREEN)], size=28)
bullets(s, 0.9, 2.1, 5.7, 3.0, "Next, with Benjamin and the GI team", PURPLE, [
    "What the Genomic Intelligence models are and where they fit.",
    "Live scenarios: GI skills running inside ClawBio on real bio problems.",
    "Follow along in Codex; prompts and links will be in the chat.",
])
bullets(s, 6.85, 2.1, 5.7, 3.0, "Try it yourself, now", GREEN, [
    "Ask a real, openly published genome questions, live: every answer executed by a ClawBio skill.",
    ("conversational.clawbio.ai", ACCENT),
    "Star the repo: github.com/ClawBio/ClawBio",
])
caption(s, 5.5, [("The point of the demos: not “the agent can run it”, but ", TEXT),
                 ("“you can trust what it ran”", GREEN), (".", TEXT)])

# ---------------------------------------------------------------- slide 12: book promo
s = slide()
chip(s, "Out now", ORANGE)
title(s, [("Build an Agentic Genomics System ", TEXT), ("(From Scratch)", MUTED)], size=26)
textbox(s, 1.4, 1.75, SW - 2.8, 0.8,
        [{"runs": [("A hands-on, build-it-yourself guide. You will not just read about agentic genomics, "
                    "you will build it in your browser on a real, openly published genome with ClawBio.", MUTED)],
          "size": 14, "align": PP_ALIGN.CENTER}])
picture(s, os.path.join(FIG, "book-qr.png"), 2.5, 2.5, cx=3.6, top=2.95)
textbox(s, 2.0, 5.5, 3.2, 1.1, [
    {"runs": [("Scan to get the book", ORANGE, True)], "size": 14, "align": PP_ALIGN.CENTER},
    {"runs": [("Paperback & Kindle  ·  328 pages", MUTED)], "size": 11.5, "align": PP_ALIGN.CENTER},
    {"runs": [("amazon.co.uk/dp/B0H6GN1ZVZ", ACCENT)], "size": 11, "align": PP_ALIGN.CENTER},
])
bullets(s, 6.8, 3.15, 5.7, 2.8, "What you will build", ORANGE, [
    "Your first runnable skill, in minutes, no setup required.",
    "Agents that choose a method, run it, read the result, and notice when it is wrong.",
    "Validation and guardrails to keep them honest.",
])

# ---------------------------------------------------------------- slide 13: community
s = slide()
chip(s, "Keep in touch")
title(s, [("Join the ", TEXT), ("ClawBio", GREEN), (" community", TEXT)], size=28)
picture(s, os.path.join(FIG, "whatsapp-qr.png"), 2.7, 2.7, cx=4.0, top=2.0)
textbox(s, 2.5, 4.8, 3.0, 0.8, [
    {"runs": [("Join the WhatsApp group", GREEN, True)], "size": 14, "align": PP_ALIGN.CENTER},
    {"runs": [("Scan with the WhatsApp camera", DIM)], "size": 11, "align": PP_ALIGN.CENTER},
])
picture(s, os.path.join(FIG, "luma-qr.png"), 2.3, 2.3, cx=9.2, top=2.1)
textbox(s, 7.7, 4.55, 3.0, 1.0, [
    {"runs": [("Subscribe to events", ACCENT, True)], "size": 14, "align": PP_ALIGN.CENTER},
    {"runs": [("Hackathons & workshops", DIM)], "size": 11, "align": PP_ALIGN.CENTER},
    {"runs": [("luma.com/ClawBio", ACCENT)], "size": 11, "align": PP_ALIGN.CENTER},
])
textbox(s, 1.0, 6.1, SW - 2.0, 0.4,
        [{"runs": [("Try it live: conversational.clawbio.ai   ·   Build with us: github.com/ClawBio/ClawBio", MUTED)],
          "size": 12.5, "align": PP_ALIGN.CENTER}])
textbox(s, 1.5, 6.55, SW - 3.0, 0.6,
        [{"runs": [("Open call: ship your own use case on the platform over the next two weeks. Best submissions win. "
                    "The book and the pharmacogenomics preprint are linked on the previous slides.", DIM)],
          "size": 10.5, "align": PP_ALIGN.CENTER}])

prs.save(OUT)
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
